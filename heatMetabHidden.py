############################################################################
#example defaults for practice (delete later, start here)
command_line = False
import os
os.chdir("C:\\Users\\Jamie\\Dropbox (Personal)\\hotMetab")
class Test:
    pass
args = Test()
args.paramfile = "C:\\Users\\Jamie\\Dropbox (Personal)\\hotMetab\\paramfile.json"
args.infile = "C:\\Users\\Jamie\\Dropbox (Personal)\\hotMetab\\inputfile.txt"

############################################################################
#### load modules ####
import os
import json
import argparse
import hot
import graph_force
import random
import datetime #datetime.date.today().strftime("%Y%m%d") will give "20230904" for example, for file names
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE
from pptx.enum.text import PP_ALIGN

############################################################################
#### additional parameters I don't think users want to set, but could move to json if users want control ####
MAX_FONT_SIZE = 12
LINE_WIDTH_PT = 3
PPT_WIDTH_IN = 13.333
PPT_HEIGHT_IN = 7.5
BOX_WIDTH_IN = 1.0
BOX_HEIGHT_IN = 0.25
SLD_LAYOUT_TITLE_AND_CONTENT = 6 #6 for blank slide layout (beware if someone has customized there Powerpoint, may need to pick a different number)
SVG_PREFIX = 'colorbar'
### global variables
hiddenNode = hot.Metabolite('hidden_node',0,1)
hiddenNode.edgeCount = 1
hiddenEdgeType = 'hidden_edge'

############################################################################
#### parse command line arguments ####
if command_line:
    parser = argparse.ArgumentParser(
        prog='heatMetab',
        description='Visualization for metabolomics data by fold-change and FDR values',
        epilog='Ask Jamie if you need help'
    )
    parser.add_argument('-i', '--infile') #input data file
    parser.add_argument('-p', '--paramfile') #parameters file (json format)
    args = parser.parse_args()

############################################################################
#### read inputs ####

# process paramfile json
paramDict = {}
print('Processing',args.paramfile)
with open(args.paramfile, 'r', encoding="utf-8") as file:
    paramDict = json.load(file)
# unfortunately, json doesn't allow numers (int, float) as keys to dicts, so I need to now convert
paramDict['valuesToColors'] = {float(k):v for k,v in paramDict['valuesToColors'].items()}
# FDRthreshold should be good as float as it is not a key, but just in case the user typed it as a string:
paramDict['FDRthreshold'] = float(paramDict['FDRthreshold'])

fc_floor = min(paramDict['valuesToColors'].keys())
fc_ceil = max(paramDict['valuesToColors'].keys())
print('Done processing parameters', paramDict)
print('The fold changes will be floored and ceilinged to',fc_floor,'and',fc_ceil)

# process infile
# infile is tab-delimited (.txt, .tab):
##column 1: metabolite names
##column 2: log-transformed fold-change values
##column 3: q-values
metabDict = {hiddenNode.name:hiddenNode} #name:Metabolite
print('Processing',args.infile)
with open(args.infile, 'r', encoding="utf-8") as file:
    for my_line in file:
        split_line = my_line.rstrip().split('\t') #name, fc, q
        split_line[1] = float(split_line[1])
        split_line[2] = float(split_line[2])
        #add to the Metabolite class set
        metabDict[split_line[0]] = hot.Metabolite(split_line[0], split_line[1], split_line[2])
fc_max = max(o.fc for o in metabDict.values())
fc_min = min(o.fc for o in metabDict.values())
print('Found',len(metabDict),'metabolites with min fold-change',fc_min,'and max fold-change',fc_max)

# process edgefile
# edgefile is tab-delimited (.txt, .tab) and contains edges, along with edge type in the 3rd column (- for normal, --- for indirect dotted line, + for combined/split pairs)
# Nicotinate	NaMN	-
# NaMN	NAD	-
# NaMN	Quinolinate	-
# Quinolinate	Tryptophan	---
# NAD	NMN	-
# NAD	NADH	-
# NAD	NaMN	-
# Glyceraldehyde 3-phosphate	Erythrose 4-phosphate	-
# Glyceraldehyde 3-phosphate	Fructose 6-phosphate	-
# Fructose 6-phosphate	Erythrose 4-phosphate	+
# ...
edgeSet = set()
lonersSet = set()
print('Processing',paramDict['edgefile'])
with open(paramDict['edgefile'], 'r', encoding="utf-8") as file:
    for my_line in file:
        split_line = my_line.rstrip().split('\t') #node1 name, node2 name, edgetype
        for mymetab in split_line[0], split_line[1]:
            if mymetab not in metabDict: #handle metabolites in the pathways that were not profiled by adding them with logFC = 0 and FDR = 1
                metabDict[mymetab] = hot.Metabolite(mymetab, 0, 1)
        node1 = metabDict[split_line[0]].name
        node2 = metabDict[split_line[1]].name
        metabDict[node1].edgeCount += 1
        metabDict[node2].edgeCount += 1
        edgeSet.add(hot.Edge(node1,node2,split_line[2]))
print('Total of',len(edgeSet),'edges')
# renumber the node indices of the metabolites such that only connected (edgeCount > 0) metabolites have a node index (starting at 0 and increasing by 1 for each) and others have node index of -1
nodeIndexCounter = 0 #will also be used later in the number of nodes for graph_force
for m in metabDict.values():
    if m.edgeCount < 1:
        m.nodeIndex = -1
        lonersSet.add(m)
        print(m.name,'has no connections to other metabolites')
    else:
        m.nodeIndex = nodeIndexCounter
        #connect each node additionally to a central hidden node for plotting
        if m != hiddenNode:
            edgeSet.add(hot.Edge(m.name,hiddenNode.name,hiddenEdgeType))
        #node counter increment
        nodeIndexCounter += 1
#make a list of tuples of node indices for graph_force
edgeList = [(metabDict[x.edgeTuple[0]].nodeIndex,metabDict[x.edgeTuple[1]].nodeIndex) for x in edgeSet]

# process namefile
# namefile is tab-delimited (.txt, .tab) containing 1st column = full name (used in the infile and edgefile) and 2nd column = display name (no first-row header)
print('Processing',paramDict['namefile'])
with open(paramDict['namefile'], 'r', encoding="utf-8") as file:
    for my_line in file:
        split_line = my_line.rstrip().split('\t') #fullname, displayname
        if split_line[0] in metabDict:
            metabDict[split_line[0]].displayName = split_line[1]
        else:
            print(split_line[0],'is not a full name in the edgefile or the infile')
print('Done processing the full name - display name file')

############################################################################
### make the color mapping ###

mycolormap = hot.colorMapMaker(valueToColorDict=paramDict['valuesToColors'],svgPrefix=SVG_PREFIX)

for m in metabDict.values():
    m.gimmeColors(colormap=mycolormap, dataFloor=fc_floor, dataCeil=fc_ceil, fdrThresh=paramDict['FDRthreshold'])

############################################################################
### make a force-directed graph ###
#use graph_force to return a list of xy-tuples coordinates list of the node locations (which is in order)
#for graph_force, node indices start at 0
#model could be creator's custom "spring_model" or "networkx_model", initial_pos could be set to coords but None is random start
#coordList = graph_force.layout_from_edge_list(number_of_nodes=len(metabDict), edges=edgeList, iter=500, model="networkx_model", initial_pos=None)
coordList = graph_force.layout_from_edge_list(number_of_nodes=nodeIndexCounter, edges=edgeList, iter=500, model="networkx_model", initial_pos=None)

#normalize the coordinate values to fit on a Powerpoint slide of widescreen size 13.333 x 7.5 in
xmin = min([t[0] for t in coordList])
xmax = max([t[0] for t in coordList])
ymin = min([t[1] for t in coordList])
ymax = max([t[1] for t in coordList])
normalizedCoordList = [((t[0]-xmin)/(xmax-xmin)*(PPT_WIDTH_IN-BOX_WIDTH_IN),(t[1]-ymin)/(ymax-ymin)*(PPT_HEIGHT_IN-BOX_HEIGHT_IN)) for t in coordList]

#do this for the lonersSet too just for even spacing
#really should make this a separate normalizeTheseCoordinates function
coordListLoners = graph_force.layout_from_edge_list(len(lonersSet), edges=[], iter=500, model="networkx_model", initial_pos=None)
xminLoners = min([t[0] for t in coordListLoners])
xmaxLoners = max([t[0] for t in coordListLoners])
yminLoners = min([t[1] for t in coordListLoners])
ymaxLoners = max([t[1] for t in coordListLoners])
normalizedCoordListLoners = [((t[0]-xminLoners)/(xmaxLoners-xminLoners)*(PPT_WIDTH_IN-BOX_WIDTH_IN),(t[1]-yminLoners)/(ymaxLoners-yminLoners)*(PPT_HEIGHT_IN-BOX_HEIGHT_IN)) for t in coordListLoners]

############################################################################
### plot to Powerpoint ###

#make the output file name
my_timestamp = str(datetime.date.today().year) + str(datetime.date.today().month) + str(datetime.date.today().day) + str(datetime.datetime.now().hour) + str(datetime.datetime.now().minute) + str(datetime.datetime.now().second)
outputName = my_timestamp + 'hotMetabOutput.pptx'

#make the presentation, slide, and shapes tree
prs = Presentation()
prs.slide_width = Inches(PPT_WIDTH_IN)
prs.slide_height = Inches(PPT_HEIGHT_IN)
slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT])
shapes = slide.shapes #_BaseShape shapes on this slide, which can contain text

# https://python-pptx.readthedocs.io/en/latest/api/shapes.html
# https://python-pptx.readthedocs.io/en/latest/api/enum/MsoLineDashStyle.html for dash styles
# later for arrowheads https://stackoverflow.com/questions/58792955/changing-format-of-connector-to-an-arrow-one-in-python-pptx
for myedge in edgeSet:
    if myedge.edgeType == hiddenEdgeType:
        continue
    (begin_x, begin_y) = normalizedCoordList[metabDict[myedge.edgeTuple[0]].nodeIndex]
    begin_x_centered = begin_x + BOX_WIDTH_IN/2
    begin_y_centered = begin_y + BOX_HEIGHT_IN/2
    (end_x, end_y) = normalizedCoordList[metabDict[myedge.edgeTuple[1]].nodeIndex]
    end_x_centered = end_x + BOX_WIDTH_IN/2
    end_y_centered = end_y + BOX_HEIGHT_IN/2
    if myedge.edgeType == '+':
        left = Inches( (begin_x+end_x)/2 + BOX_WIDTH_IN/2 - BOX_WIDTH_IN/8)
        top = Inches( (begin_y+end_y)/2)
        width = Inches(BOX_WIDTH_IN/4)
        height = width
        shape = shapes.add_shape(MSO_SHAPE.MATH_PLUS, left, top, width, height)
        shape.shadow.inherit = False #turn off shadow
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(128,128,128)
        line = shape.line
        line.color.rgb = RGBColor(255,255,255)
        line.width = Pt(0)
        #add a gray connector to ensure the + is located in the right place, user could delete
        connector = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(begin_x_centered), Inches(begin_y_centered), Inches(end_x_centered), Inches(end_y_centered))
        if myedge.edgeType == '---':
            connector.line.dash_style = MSO_LINE.DASH
        connector.shadow.inherit = False
        connector.line.fill.background()
        connector.line.fill.solid()
        connector.line.width = Pt(1)
        connector.line.fill.fore_color.rgb = RGBColor(128,128,128)
    else:
        connector = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(begin_x_centered), Inches(begin_y_centered), Inches(end_x_centered), Inches(end_y_centered))
        if myedge.edgeType == '---':
            connector.line.dash_style = MSO_LINE.DASH
        connector.shadow.inherit = False
        connector.line.fill.background()
        connector.line.fill.solid()
        connector.line.width = Pt(1)
        connector.line.fill.fore_color.rgb = RGBColor(0,0,0)

#add the boxed words
for m in metabDict.values():
    if m == hiddenNode:
        continue
    #position and size
    n = m.nodeIndex
    if n < 0: # unconnected node, plot on next slide instead
        continue
    left = Inches(normalizedCoordList[n][0]) #x
    top = Inches(normalizedCoordList[n][1]) #y
    width = Inches(BOX_WIDTH_IN)
    height = Inches(BOX_HEIGHT_IN)
    #fill and outline color
    myfillcolor = tuple(int(x*255) for x in m.fillCol[0:3]) #convert from 0-1 to 0-255 scale and remove the alpha
    myoutlinecolor = tuple(int(x*255) for x in m.outlineCol[0:3]) #convert from 0-1 to 0-255 scale and remove the alpha
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height) #ROUNDED_RECTANGLE is another shape option
    shape.shadow.inherit = False #turn off shadow
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*myfillcolor)
    line = shape.line
    line.color.rgb = RGBColor(*myoutlinecolor)
    line.width = Pt(LINE_WIDTH_PT)
    #text
    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = str(m.displayName)
    font = run.font
    font.name = 'Arial Narrow'
    font.color.rgb = RGBColor(0,0,0)
    font.size = Pt(8)
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.word_wrap = True
    #text_frame.fit_text(max_size=4) #errors out if text is too large, so omit; otherwise would autofit text
    #text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE #part of the above

#singletons on separate slide
#would be better to make a plottingMetabolites function
slide_loners = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT])
for count,m in enumerate(lonersSet):
        left = Inches(normalizedCoordListLoners[count][0]) #x
        top = Inches(normalizedCoordListLoners[count][1]) #y
        width = Inches(BOX_WIDTH_IN)
        height = Inches(BOX_HEIGHT_IN)
        #fill and outline color
        myfillcolor = tuple(int(x*255) for x in m.fillCol[0:3]) #convert from 0-1 to 0-255 scale and remove the alpha
        myoutlinecolor = tuple(int(x*255) for x in m.outlineCol[0:3]) #convert from 0-1 to 0-255 scale and remove the alpha
        shape = slide_loners.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height) #ROUNDED_RECTANGLE is another shape option
        shape.shadow.inherit = False #turn off shadow
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*myfillcolor)
        line = shape.line
        line.color.rgb = RGBColor(*myoutlinecolor)
        line.width = Pt(LINE_WIDTH_PT)
        #text
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = str(m.displayName)
        font = run.font
        font.name = 'Arial Narrow'
        font.color.rgb = RGBColor(0,0,0)
        font.size = Pt(8)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.word_wrap = True
txBox_loners = slide_loners.shapes.add_textbox(Inches(PPT_WIDTH_IN/4),0,Inches(PPT_WIDTH_IN/2),Inches(2))
tf_loners = txBox_loners.text_frame
tf_loners.text = "Metabolites detected without pathway connections in the edge file"

#colorbar legend on new slide; pdf and svg versions of the colorbar are also saved in the cwd
slide_leg = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT])
txBox_leg = slide_leg.shapes.add_textbox(Inches(PPT_WIDTH_IN/4),0,Inches(PPT_WIDTH_IN/2),Inches(2))
tf_leg = txBox_leg.text_frame
tf_leg.text = "Legend"
slide_leg.shapes.add_picture(SVG_PREFIX + '.png',left=Inches(PPT_WIDTH_IN/4),top=Inches(PPT_HEIGHT_IN/4))

#save
prs.save(outputName)

############################################################################
### all done ###
print('All done! Check this current folder or elsewhere if you specified for your output files',os.getcwd())

###connect all subgraphs to a (?hidden) root node for best force directed graph appearance (not too pushed to the side)
###or try d3
###or graph_force separately each subgraph and plot it on a fraction of a slide
