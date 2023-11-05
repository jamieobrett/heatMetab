############################################################################
#example defaults for practice (use command_line = False if using jupyterlab)
command_line = True
if not command_line:
    import os
    os.chdir("C:\\Users\\Jamie\\Dropbox (Personal)\\hotMetab")
    class Test:
        pass
    args = Test()
    args.paramfile = "C:\\Users\\Jamie\\Dropbox (Personal)\\hotMetab\\paramfile.json"
    args.infile = "C:\\Users\\Jamie\\Dropbox (Personal)\\hotMetab\\inputfile.txt"

############################################################################
#### load modules ####
import os
import json
import argparse
import hot
import random
import datetime #datetime.date.today().strftime("%Y%m%d") will give "20230904" for example, for file names
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE
from pptx.enum.text import PP_ALIGN

############################################################################
#### additional parameters I don't think users want to set, but could move to json if users want control ####
MAX_FONT_SIZE = 12
LINE_WIDTH_PT = 3
PPT_WIDTH_IN = 13.333
PPT_HEIGHT_IN = 7.5
BOX_WIDTH_IN = 0.5
BOX_HEIGHT_IN = 0.15
FONT_SIZE_PT = 6
SLD_LAYOUT_TITLE_AND_CONTENT = 6 #6 for blank slide layout (beware if someone has customized there Powerpoint, may need to pick a different number)
COLORBARIMAGE_PREFIX = 'colorbar'

############################################################################
#### parse command line arguments ####
if command_line:
    parser = argparse.ArgumentParser(
        prog='heatMetab',
        description='Visualization for metabolomics data by fold-change and FDR values',
        epilog='Ask Jamie if you need help'
    )
    parser.add_argument('-i', '--infile') #input data file
    parser.add_argument('-p', '--paramfile') #parameters file (json format)
    args = parser.parse_args()

############################################################################
#### read inputs ####
paramDict = {} #populated by the json paramfile
metabDict = {} #dict of name:Metabolite of all metabolites detected plus all metabolites plotted, which allows lookup for coloring
edgeSet = set() #set of Edge instances, which can be iterated over for plotting
nodeList = [] #list of Node instances, importantly ordered by node index, which allows lookup for Edge ends
lonersSet = set() #set of Metabolite instances to be plotted in a separate slide
lonersCoords = [] #a list of (x,y) coordinates for neatly plotting the loners (whose arrangement otherwise doesn't matter)

# process paramfile json
print('Processing',args.paramfile)
with open(args.paramfile, 'r', encoding="utf-8") as file:
    paramDict = json.load(file)
# unfortunately, json doesn't allow numers (int, float) as keys to dicts, so I need to now convert
paramDict['valuesToColors'] = {float(k):v for k,v in paramDict['valuesToColors'].items()}
# FDRthreshold should be good as float as it is not a key, but just in case the user typed it as a string:
paramDict['FDRthreshold'] = float(paramDict['FDRthreshold'])

fc_floor = min(paramDict['valuesToColors'].keys())
fc_ceil = max(paramDict['valuesToColors'].keys())
print('Done processing parameters', paramDict)
print('The fold changes will be floored and ceilinged to',fc_floor,'and',fc_ceil)

# process infile
# tab-delimited (.txt, .tab):
##column 1: metabolite names
##column 2: log-transformed fold-change values
##column 3: q-values
print('Processing',args.infile)
with open(args.infile, 'r', encoding="utf-8") as file:
    for my_line in file:
        split_line = my_line.rstrip().split('\t') #name, fc, q
        split_line[1] = float(split_line[1])
        split_line[2] = float(split_line[2])
        #add to the Metabolite class set
        metabDict[split_line[0]] = hot.Metabolite(split_line[0], split_line[1], split_line[2])
fc_max = max(o.fc for o in metabDict.values())
fc_min = min(o.fc for o in metabDict.values())
print('Found',len(metabDict),'metabolites with min fold-change',fc_min,'and max fold-change',fc_max)

# process coordfile
# tab-delimited (.txt, .tab)
## node index number
## node (metabolite) name
## xcoord (in EMU) of the textbox center
## ycoord (in EMU) of the textbox center
## number of connections
## comma-delimited list of connected node index numbers
print('Processing',paramDict['coordfile'])
with open(paramDict['coordfile'], 'r', encoding="utf-8") as file:
    for my_line in file:
        (node_index, node_name, x_center, y_center, n_connections) = my_line.rstrip().split('\t')[0:5]
        if len(my_line.rstrip().split('\t')) == 6:
            connections_list = my_line.rstrip().split('\t')[5]
        else:
            connections_list = None
        if node_name not in metabDict: #handle metabolites in the pathways that were not profiled by adding them with logFC = None and FDR = None
            metabDict[node_name] = hot.Metabolite(node_name,None,None)
        metabDict[node_name].edgeCount = int(n_connections)
        metabDict[node_name].nodeIndices.append(int(node_index))
        nodeList.append(hot.Node(node_name,int(node_index),float(x_center),float(y_center)))
        if connections_list:
            connections_list_aslist = [int(c) for c in connections_list.split(',')]
            for connection_node_index in connections_list_aslist:
                edgeSet.add(hot.Edge(int(node_index), connection_node_index))
print('Total of',len(edgeSet),'edges')

# process namefile
# namefile is tab-delimited (.txt, .tab) containing 1st column = full name (used in the infile and edgefile) and 2nd column = display name (no first-row header)
print('Processing',paramDict['namefile'])
with open(paramDict['namefile'], 'r', encoding="utf-8") as file:
    for my_line in file:
        split_line = my_line.rstrip().split('\t') #fullname, displayname
        if split_line[0] in metabDict:
            metabDict[split_line[0]].displayName = split_line[1]
        else:
            print(split_line[0],'is not a full name in the edgefile or the infile')
print('Done processing the full name - display name file')

############################################################################
### make the color mapping ###

mycolormap = hot.colorMapMaker(valueToColorDict=paramDict['valuesToColors'],colorbarImagePrefix=COLORBARIMAGE_PREFIX)

for m in metabDict.values():
    m.gimmeColors(colormap=mycolormap, dataFloor=fc_floor, dataCeil=fc_ceil, fdrThresh=paramDict['FDRthreshold'], undetectedColor=paramDict['undetectedColor'])

############################################################################
### normalize the coordinates to fit and handle loners ###

coordList = [(n.center_x,n.center_y) for n in nodeList]
normalizedCoordList = hot.normalizeTheseCoordinates(coordList=coordList,boxwidth=PPT_WIDTH_IN-BOX_WIDTH_IN,boxheight=PPT_HEIGHT_IN-BOX_HEIGHT_IN,offsetx=BOX_WIDTH_IN,offsety=BOX_HEIGHT_IN)
for count,n in enumerate(nodeList):
    n.center_x = normalizedCoordList[count][0]
    n.center_y = normalizedCoordList[count][1]

# make the lonersSet
for m in metabDict.values():
    if m.edgeCount < 1:
        lonersSet.add(m)
        print(m.name,'has no connections to other metabolites')

# make coordinates for lonersSet
running_x_loners = BOX_WIDTH_IN/2
running_y_loners = BOX_HEIGHT_IN * 4
for loner in lonersSet:
    lonersCoords.append((running_x_loners, running_y_loners))
    if running_x_loners < (PPT_WIDTH_IN - BOX_WIDTH_IN*2):
        running_x_loners += BOX_WIDTH_IN*1.1
    else:
        running_x_loners = BOX_WIDTH_IN/2
        running_y_loners += BOX_HEIGHT_IN*2

############################################################################
### plot to Powerpoint ###

# make the output file name
my_timestamp = str(datetime.date.today().year) + str(datetime.date.today().month) + str(datetime.date.today().day) + str(datetime.datetime.now().hour) + str(datetime.datetime.now().minute) + str(datetime.datetime.now().second)
outputName = my_timestamp + 'hotMetabOutput.pptx'

# make the presentation, slide, and shapes tree
prs = Presentation()
prs.slide_width = Inches(PPT_WIDTH_IN)
prs.slide_height = Inches(PPT_HEIGHT_IN)
slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT])
shapes = slide.shapes #_BaseShape shapes on this slide, which can contain text

# edges first
# https://python-pptx.readthedocs.io/en/latest/api/shapes.html
# https://python-pptx.readthedocs.io/en/latest/api/enum/MsoLineDashStyle.html for dash styles
# later for arrowheads https://stackoverflow.com/questions/58792955/changing-format-of-connector-to-an-arrow-one-in-python-pptx
for myedge in edgeSet:
    begin_x = nodeList[myedge.edgeTuple[0]].center_x
    begin_y = nodeList[myedge.edgeTuple[0]].center_y
    end_x = nodeList[myedge.edgeTuple[1]].center_x
    end_y = nodeList[myedge.edgeTuple[1]].center_y
    connector = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(begin_x), Inches(begin_y), Inches(end_x), Inches(end_y))
    connector.shadow.inherit = False
    connector.line.fill.background()
    connector.line.fill.solid()
    connector.line.width = Pt(1)
    connector.line.fill.fore_color.rgb = RGBColor(0,0,0)

# add the boxed words
for n in nodeList:
    if metabDict[n.name].edgeCount < 1: # skip loners
        continue
    else:
        m = metabDict[n.name]
        left = Inches(n.center_x - BOX_WIDTH_IN/2)
        top = Inches(n.center_y - BOX_HEIGHT_IN/2)
        width = Inches(BOX_WIDTH_IN)
        height = Inches(BOX_HEIGHT_IN)
        #fill and outline color
        myfillcolor = tuple(int(x*255) for x in m.fillCol[0:3]) #convert from 0-1 to 0-255 scale and remove the alpha
        myoutlinecolor = tuple(int(x*255) for x in m.outlineCol[0:3]) #convert from 0-1 to 0-255 scale and remove the alpha
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height) #ROUNDED_RECTANGLE is another shape option
        shape.shadow.inherit = False #turn off shadow
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*myfillcolor)
        line = shape.line
        line.color.rgb = RGBColor(*myoutlinecolor)
        line.width = Pt(LINE_WIDTH_PT)
        #text
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = str(m.displayName)
        font = run.font
        font.name = 'Arial Narrow'
        font.color.rgb = RGBColor(0,0,0)
        font.size = Pt(FONT_SIZE_PT)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.word_wrap = True
        #text_frame.fit_text(max_size=4) #errors out if text is too large, so omit; otherwise would autofit text
        #text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE #part of the above

#singletons on separate slide
slide_loners = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT])
for count,m in enumerate(lonersSet):
        left = Inches(lonersCoords[count][0]) #x
        top = Inches(lonersCoords[count][1]) #y
        width = Inches(BOX_WIDTH_IN)
        height = Inches(BOX_HEIGHT_IN)
        #fill and outline color
        myfillcolor = tuple(int(x*255) for x in m.fillCol[0:3]) #convert from 0-1 to 0-255 scale and remove the alpha
        myoutlinecolor = tuple(int(x*255) for x in m.outlineCol[0:3]) #convert from 0-1 to 0-255 scale and remove the alpha
        shape = slide_loners.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height) #ROUNDED_RECTANGLE is another shape option
        shape.shadow.inherit = False #turn off shadow
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*myfillcolor)
        line = shape.line
        line.color.rgb = RGBColor(*myoutlinecolor)
        line.width = Pt(LINE_WIDTH_PT)
        #text
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = str(m.displayName)
        font = run.font
        font.name = 'Arial Narrow'
        font.color.rgb = RGBColor(0,0,0)
        font.size = Pt(FONT_SIZE_PT)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.word_wrap = True
txBox_loners = slide_loners.shapes.add_textbox(Inches(PPT_WIDTH_IN/4),0,Inches(PPT_WIDTH_IN/2),Inches(2))
tf_loners = txBox_loners.text_frame
tf_loners.text = "Metabolites detected without any connections to other metabolites"

#colorbar legend on new slide; pdf and svg versions of the colorbar are also saved in the cwd
slide_leg = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT])
txBox_leg = slide_leg.shapes.add_textbox(Inches(PPT_WIDTH_IN/4),0,Inches(PPT_WIDTH_IN/2),Inches(2))
tf_leg = txBox_leg.text_frame
tf_leg.text = "Legend. Color shows change magnitude. Fill is present for FDR " + str(paramDict['FDRthreshold'])
slide_leg.shapes.add_picture(COLORBARIMAGE_PREFIX + '.png',left=Inches(PPT_WIDTH_IN/4),top=Inches(PPT_HEIGHT_IN/4))
#add the 'not detected' legend
left = Inches(PPT_WIDTH_IN/2)
top = Inches(PPT_HEIGHT_IN/2)
width = Inches(BOX_WIDTH_IN)
height = Inches(BOX_HEIGHT_IN)
undetectedDummy = hot.Metabolite("undetectedDummy",None,None)
undetectedDummy.gimmeColors(colormap=mycolormap, dataFloor=fc_floor, dataCeil=fc_ceil, fdrThresh=paramDict['FDRthreshold'], undetectedColor=paramDict['undetectedColor'])
myfillcolor = (255,255,255)
myoutlinecolor = tuple(int(x*255) for x in undetectedDummy.outlineCol[0:3])
shape = slide_leg.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
shape.shadow.inherit = False
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(*myfillcolor)
line = shape.line
line.color.rgb = RGBColor(*myoutlinecolor)
line.width = Pt(LINE_WIDTH_PT)
text_frame = shape.text_frame
p = text_frame.paragraphs[0]
run = p.add_run()
run.text = str('Not detected')
font = run.font
font.name = 'Arial Narrow'
font.color.rgb = RGBColor(0,0,0)
font.size = Pt(FONT_SIZE_PT)
text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
text_frame.word_wrap = True

#save
prs.save(outputName)

############################################################################
### all done ###
print('All done! Check this current folder or elsewhere if you specified for your output files',os.getcwd())
