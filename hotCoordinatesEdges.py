# Extract the coordinates from a Powerpoint slide of metabolites and edges in pathways made manually by a user
# Outputs a tab-delimited file containing these columns (no header):
## node index number
## node (metabolite) name; note this is whatever is in the Powerpoint file and should be converted to match the edge file names and input file names
## xcoord (in EMU) of the textbox center
## ycoord (in EMU) of the textbox center
## number of connections
## comma-delimited list of connected node index numbers
# later I could change this to have + and --- connectors

#example defaults for practice (use command_line = False if using jupyterlab)
command_line = True
if not command_line:
    import os
    os.chdir("C:\\Users\\Jamie\\Dropbox (Personal)\\hotMetab")
    class Test:
        pass
    args = Test()
    args.infile = "C:\\Users\\Jamie\\Dropbox (Personal)\\hotMetab\\inputfile.txt"
    args.infile = "C:\\Users\\Jamie\\Dropbox (Personal)\\hotMetab\\hardcodedCoordinates.txt"

import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

############################################################################
#### parse command line arguments ####
if command_line:
    parser = argparse.ArgumentParser(
        prog='hotCoordinatesEdges',
        description='Helper for heatMetab to make the coordinates and edges file from a Powerpoint slide',
        epilog='Ask Jamie if you need help'
    )
    parser.add_argument('-i', '--infile') #input Powerpoint file name
    parser.add_argument('-o', '--outfile') #output file name
    args = parser.parse_args()

# parameters
infile = args.infile
outfile = args.outfile

# classes and functions
class MetabNode:
    """A class for metabolite nodes"""
    def __init__(self, name, left, top, width, height):
        self.nodeIndex = -1 #0-based
        self.name = name
        self.left = left
        self.top = top
        self.width = width
        self.height = height
        self.center_x = self.left + self.width/2
        self.center_y = self.top + self.height/2
        self.connections = []
        self.nconnections = 0
    def nearest_squared_distance_to_point(self, point_x, point_y):
        """Nearest distance from a point to any outline edge of this box, assuming that the distance to a point inside of this rect is zero"""
        dx = abs(self.center_x-point_x) - (self.width*0.5)
        dy = abs(self.center_y-point_y) - (self.height*0.5)
        result = (dx * (dx > 0)) ** 2 + (dy * (dy > 0)) ** 2
        return result

# open the Powerpoint
nodeDict = {} #node index to MetabNode
prs = Presentation(infile)
slide = prs.slides[0]

# get all the text boxes
for shape in slide.shapes:
    if shape.has_text_frame:
        myindex = len(nodeDict)
        nodeDict[myindex] = MetabNode(shape.text,shape.left,shape.top,shape.width,shape.height)
        nodeDict[myindex].nodeIndex = myindex

# get all the edges
for shape in slide.shapes:
    if(shape.shape_type == MSO_SHAPE_TYPE.LINE):
        #get the nearest nodes
        minimum_node1_index = -1
        minimum_node2_index = -1
        minimum_node1_dist = minimum_node2_dist = prs.slide_width ** 2 + prs.slide_height ** 2 #starting value is larger than any nearest node
        for n in nodeDict.values():
            mydist1 = n.nearest_squared_distance_to_point(shape.begin_x,shape.begin_y)
            mydist2 = n.nearest_squared_distance_to_point(shape.end_x,shape.end_y)
            if mydist1 < minimum_node1_dist:
                minimum_node1_index = n.nodeIndex
                minimum_node1_dist = mydist1
            if mydist2 < minimum_node2_dist:
                minimum_node2_index = n.nodeIndex
                minimum_node2_dist = mydist2
        nodeDict[minimum_node1_index].connections.append(minimum_node2_index)
        nodeDict[minimum_node1_index].nconnections += 1
        nodeDict[minimum_node2_index].nconnections += 1

with open (outfile, 'w', encoding='utf-8') as file:
    for m in nodeDict.values():
        outTuple = ( str(m.nodeIndex), m.name, str(m.center_x), str(m.center_y), str(m.nconnections), ','.join(str(i) for i in m.connections) )
        outLine = "\t".join(outTuple)
        file.write(outLine + "\n")